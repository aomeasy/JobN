"""
ยูทิลิตี้สำหรับจัดการไฟล์
อัพโหลด, บันทึก, สกัดข้อความ
"""

import os
import uuid
import mimetypes
import hashlib
from datetime import datetime
from typing import Optional, Dict, Any, List
import logging
from pathlib import Path

import streamlit as st
from PIL import Image
import PyPDF2
import docx
import openpyxl
from pptx import Presentation

from config import config
from database.database import get_db_session
from database.models import Document

logger = logging.getLogger(__name__)

class FileHandler:
    """คลาสจัดการไฟล์"""
    
    def __init__(self):
        self.upload_folder = Path(config.app.upload_folder)
        self.upload_folder.mkdir(parents=True, exist_ok=True)
        
    def save_uploaded_file(self, uploaded_file, category: str = None, 
                          tags: List[str] = None, is_public: bool = False,
                          user_id: int = 1) -> Optional[Dict[str, Any]]:
        """บันทึกไฟล์ที่อัพโหลดและข้อมูลลงฐานข้อมูล"""
        
        try:
            # สร้างชื่อไฟล์ใหม่
            file_uuid = str(uuid.uuid4())
            file_extension = Path(uploaded_file.name).suffix.lower()
            new_filename = f"{file_uuid}{file_extension}"
            
            # เส้นทางไฟล์
            file_path = self.upload_folder / new_filename
            
            # บันทึกไฟล์
            with open(file_path, "wb") as f:
                f.write(uploaded_file.getvalue())
            
            # ดึงข้อมูลไฟล์
            file_size = file_path.stat().st_size
            mime_type = mimetypes.guess_type(uploaded_file.name)[0] or 'application/octet-stream'
            
            # สกัดข้อความ
            extracted_text = self.extract_text_from_file(file_path, file_extension)
            
            # บันทึกข้อมูลลงฐานข้อมูล
            with get_db_session() as session:
                document = Document(
                    filename=new_filename,
                    original_filename=uploaded_file.name,
                    file_path=str(file_path),
                    file_size=file_size,
                    file_type=file_extension[1:],  # ลบจุดออก
                    mime_type=mime_type,
                    category=category or "เอกสารทั่วไป",
                    tags=tags,
                    extracted_text=extracted_text,
                    is_processed=bool(extracted_text),
                    processing_status="completed" if extracted_text else "pending",
                    uploaded_by=user_id,
                    is_public=is_public
                )
                
                session.add(document)
                session.commit()
                session.refresh(document)
                
                logger.info(f"บันทึกไฟล์ {uploaded_file.name} สำเร็จ (ID: {document.id})")
                
                return {
                    "success": True,
                    "document_id": document.id,
                    "filename": new_filename,
                    "original_filename": uploaded_file.name,
                    "file_size": file_size,
                    "extracted_text": bool(extracted_text)
                }
                
        except Exception as e:
            logger.error(f"เกิดข้อผิดพลาดในการบันทึกไฟล์ {uploaded_file.name}: {e}")
            
            # ลบไฟล์ถ้าบันทึกไม่สำเร็จ
            try:
                if file_path.exists():
                    file_path.unlink()
            except:
                pass
                
            return None
    
    def extract_text_from_file(self, file_path: Path, file_extension: str) -> Optional[str]:
        """สกัดข้อความจากไฟล์ตามประเภท"""
        
        try:
            if file_extension == '.pdf':
                return self._extract_from_pdf(file_path)
            elif file_extension == '.docx':
                return self._extract_from_docx(file_path)
            elif file_extension == '.xlsx':
                return self._extract_from_xlsx(file_path)
            elif file_extension == '.pptx':
                return self._extract_from_pptx(file_path)
            elif file_extension == '.txt':
                return self._extract_from_txt(file_path)
            elif file_extension in ['.jpg', '.jpeg', '.png', '.tiff', '.bmp']:
                # สำหรับรูปภาพ จะใช้ OCR ภายหลัง
                return None
            else:
                logger.warning(f"ไม่รองรับการสกัดข้อความจากไฟล์ประเภท {file_extension}")
                return None
                
        except Exception as e:
            logger.error(f"เกิดข้อผิดพลาดในการสกัดข้อความจาก {file_path}: {e}")
            return None
    
    def _extract_from_pdf(self, file_path: Path) -> Optional[str]:
        """สกัดข้อความจาก PDF"""
        try:
            text_content = []
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        text = page.extract_text()
                        if text.strip():
                            text_content.append(f"--- หน้า {page_num + 1} ---\n{text}")
                    except Exception as e:
                        logger.warning(f"ไม่สามารถสกัดข้อความจากหน้า {page_num + 1}: {e}")
                        continue
            
            return "\n\n".join(text_content) if text_content else None
            
        except Exception as e:
            logger.error(f"เกิดข้อผิดพลาดในการอ่าน PDF: {e}")
            return None
    
    def _extract_from_docx(self, file_path: Path) -> Optional[str]:
        """สกัดข้อความจาก Word Document"""
        try:
            doc = docx.Document(file_path)
            text_content = []
            
            # สกัดข้อความจาก paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # สกัดข้อความจากตาราง
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_text.append(" | ".join(row_text))
                
                if table_text:
                    text_content.append("\n--- ตาราง ---\n" + "\n".join(table_text))
            
            return "\n\n".join(text_content) if text_content else None
            
        except Exception as e:
            logger.error(f"เกิดข้อผิดพลาดในการอ่าน Word: {e}")
            return None
    
    def _extract_from_xlsx(self, file_path: Path) -> Optional[str]:
        """สกัดข้อความจาก Excel"""
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text_content = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = [f"--- แผ่นงาน: {sheet_name} ---"]
                
                # หาขอบเขตข้อมูล
                max_row = sheet.max_row
                max_col = sheet.max_column
                
                if max_row > 100:  # จำกัดแถวไม่เกิน 100
                    max_row = 100
                if max_col > 20:   # จำกัดคอลัมน์ไม่เกิน 20
                    max_col = 20
                
                for row in range(1, max_row + 1):
                    row_data = []
                    for col in range(1, max_col + 1):
                        cell_value = sheet.cell(row=row, column=col).value
                        if cell_value is not None:
                            row_data.append(str(cell_value))
                        else:
                            row_data.append("")
                    
                    # เพิ่มแถวที่มีข้อมูล
                    if any(cell.strip() for cell in row_data):
                        sheet_text.append(" | ".join(row_data))
                
                if len(sheet_text) > 1:  # มีข้อมูลนอกจากหัวเรื่อง
                    text_content.append("\n".join(sheet_text))
            
            workbook.close()
            return "\n\n".join(text_content) if text_content else None
            
        except Exception as e:
            logger.error(f"เกิดข้อผิดพลาดในการอ่าน Excel: {e}")
            return None
    
    def _extract_from_pptx(self, file_path: Path) -> Optional[str]:
        """สกัดข้อความจาก PowerPoint"""
        try:
            presentation = Presentation(file_path)
            text_content = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = [f"--- สไลด์ {slide_num} ---"]
                
                # สกัดข้อความจาก shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:  # มีข้อมูลนอกจากหัวเรื่อง
                    text_content.append("\n".join(slide_text))
            
            return "\n\n".join(text_content) if text_content else None
            
        except Exception as e:
            logger.error(f"เกิดข้อผิดพลาดในการอ่าน PowerPoint: {e}")
            return None
    
    def _extract_from_txt(self, file_path: Path) -> Optional[str]:
        """สกัดข้อความจากไฟล์ text"""
        try:
            # ลองหลาย encoding
            encodings = ['utf-8', 'utf-8-sig', 'cp874', 'cp1252', 'latin-1']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        content = file.read()
                        if content.strip():
                            return content
                except UnicodeDecodeError:
                    continue
            
            logger.warning(f"ไม่สามารถอ่านไฟล์ text ด้วย encoding ใดๆ ได้")
            return None
            
        except Exception as e:
            logger.error(f"เกิดข้อผิดพลาดในการอ่านไฟล์ text: {e}")
            return None
    
    def get_file_info(self, file_path: Path) -> Dict[str, Any]:
        """ดึงข้อมูลไฟล์"""
        try:
            stat = file_path.stat()
            
            return {
                "size": stat.st_size,
                "created": datetime.fromtimestamp(stat.st_ctime),
                "modified": datetime.fromtimestamp(stat.st_mtime),
                "extension": file_path.suffix.lower(),
                "mime_type": mimetypes.guess_type(str(file_path))[0]
            }
        except Exception as e:
            logger.error(f"ไม่สามารถดึงข้อมูลไฟล์ได้: {e}")
            return {}
    
    def calculate_file_hash(self, file_path: Path) -> Optional[str]:
        """คำนวณ hash ของไฟล์เพื่อตรวจสอบการซ้ำ"""
        try:
            hash_md5 = hashlib.md5()
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_md5.update(chunk)
            return hash_md5.hexdigest()
        except Exception as e:
            logger.error(f"ไม่สามารถคำนวณ hash ได้: {e}")
            return None
    
    def validate_file(self, uploaded_file) -> Dict[str, Any]:
        """ตรวจสอบไฟล์ก่อนอัพโหลด"""
        validation_result = {
            "valid": True,
            "errors": [],
            "warnings": []
        }
        
        # ตรวจสอบขนาดไฟล์
        max_size = config.app.max_file_size * 1024 * 1024  # MB to bytes
        if uploaded_file.size > max_size:
            validation_result["valid"] = False
            validation_result["errors"].append(
                f"ไฟล์ใหญ่เกินกำหนด ({uploaded_file.size / 1024 / 1024:.1f}MB > {config.app.max_file_size}MB)"
            )
        
        # ตรวจสอบประเภทไฟล์
        file_extension = Path(uploaded_file.name).suffix.lower()
        supported_extensions = ['.pdf', '.docx', '.xlsx', '.pptx', '.txt', '.jpg', '.jpeg', '.png', '.tiff', '.bmp']
        
        if file_extension not in supported_extensions:
            validation_result["valid"] = False
            validation_result["errors"].append(
                f"ไม่รองรับไฟล์ประเภท {file_extension}"
            )
        
        # ตรวจสอบชื่อไฟล์
        if len(uploaded_file.name) > 255:
            validation_result["valid"] = False
            validation_result["errors"].append("ชื่อไฟล์ยาวเกินไป (เกิน 255 ตัวอักษร)")
        
        # คำเตือนสำหรับไฟล์รูปภาพ
        if file_extension in ['.jpg', '.jpeg', '.png', '.tiff', '.bmp']:
            validation_result["warnings"].append(
                "ไฟล์รูปภาพจะต้องผ่านการประมวลผล OCR เพิ่มเติม"
            )
        
        return validation_result
    
    def cleanup_temp_files(self, older_than_hours: int = 24):
        """ทำความสะอาดไฟล์ชั่วคราว"""
        try:
            temp_folder = Path("data/temp")
            if not temp_folder.exists():
                return
            
            cutoff_time = datetime.now().timestamp() - (older_than_hours * 3600)
            deleted_count = 0
            
            for file_path in temp_folder.iterdir():
                if file_path.is_file():
                    try:
                        if file_path.stat().st_mtime < cutoff_time:
                            file_path.unlink()
                            deleted_count += 1
                    except Exception as e:
                        logger.warning(f"ไม่สามารถลบไฟล์ {file_path}: {e}")
            
            logger.info(f"ลบไฟล์ชั่วคราว {deleted_count} ไฟล์")
            
        except Exception as e:
            logger.error(f"เกิดข้อผิดพลาดในการทำความสะอาด: {e}")
    
    def get_storage_usage(self) -> Dict[str, Any]:
        """ดึงข้อมูลการใช้พื้นที่"""
        try:
            total_size = 0
            file_count = 0
            
            for file_path in self.upload_folder.rglob("*"):
                if file_path.is_file():
                    total_size += file_path.stat().st_size
                    file_count += 1
            
            return {
                "total_size_bytes": total_size,
                "total_size_mb": round(total_size / 1024 / 1024, 2),
                "file_count": file_count,
                "folder_path": str(self.upload_folder)
            }
            
        except Exception as e:
            logger.error(f"ไม่สามารถดึงข้อมูลการใช้พื้นที่ได้: {e}")
            return {}

# Utility functions สำหรับใช้ใน Streamlit
def format_file_size(size_bytes: int) -> str:
    """แปลงขนาดไฟล์เป็นรูปแบบที่อ่านง่าย"""
    if size_bytes == 0:
        return "0B"
    
    import math
    size_names = ["B", "KB", "MB", "GB"]
    i = int(math.floor(math.log(size_bytes, 1024)))
    p = math.pow(1024, i)
    s = round(size_bytes / p, 2)
    return f"{s} {size_names[i]}"

def get_file_icon(file_extension: str) -> str:
    """ส่งคืนไอคอนตามประเภทไฟล์"""
    icons = {
        '.pdf': '📄',
        '.docx': '📝', 
        '.doc': '📝',
        '.xlsx': '📊',
        '.xls': '📊',
        '.pptx': '📋',
        '.ppt': '📋',
        '.txt': '📃',
        '.jpg': '🖼️',
        '.jpeg': '🖼️',
        '.png': '🖼️',
        '.gif': '🖼️',
        '.tiff': '🖼️',
        '.bmp': '🖼️'
    }
    return icons.get(file_extension.lower(), '📄')

def is_image_file(filename: str) -> bool:
    """ตรวจสอบว่าเป็นไฟล์รูปภาพหรือไม่"""
    image_extensions = ['.jpg', '.jpeg', '.png', '.gif', '.tiff', '.bmp', '.webp']
    extension = Path(filename).suffix.lower()
    return extension in image_extensions

def is_document_file(filename: str) -> bool:
    """ตรวจสอบว่าเป็นไฟล์เอกสารหรือไม่"""
    doc_extensions = ['.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.txt']
    extension = Path(filename).suffix.lower()
    return extension in doc_extensions

# สร้าง instance หลัก
file_handler = FileHandler()

# ตัวอย่างการใช้งานใน Streamlit
def show_file_upload_widget():
    """แสดง widget สำหรับอัพโหลดไฟล์พร้อมการตรวจสอบ"""
    
    uploaded_file = st.file_uploader(
        "เลือกไฟล์",
        type=['pdf', 'docx', 'xlsx', 'pptx', 'txt', 'jpg', 'jpeg', 'png'],
        help="รองรับไฟล์เอกสารและรูปภาพ"
    )
    
    if uploaded_file:
        # ตรวจสอบไฟล์
        validation = file_handler.validate_file(uploaded_file)
        
        if not validation["valid"]:
            for error in validation["errors"]:
                st.error(f"❌ {error}")
            return None
        
        if validation["warnings"]:
            for warning in validation["warnings"]:
                st.warning(f"⚠️ {warning}")
        
        # แสดงข้อมูลไฟล์
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.metric("ขนาดไฟล์", format_file_size(uploaded_file.size))
        
        with col2:
            st.metric("ประเภท", uploaded_file.type or "ไม่ทราบ")
        
        with col3:
            icon = get_file_icon(Path(uploaded_file.name).suffix)
            st.metric("รูปแบบ", f"{icon} {Path(uploaded_file.name).suffix}")
        
        return uploaded_file
    
    return None
